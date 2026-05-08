"""
REQ-006 doc_parser — adapters/parsers/pptx_parser.py

PptxParser
python-pptx 기반 PowerPoint 파서

처리 항목:
    - 슬라이드별 텍스트 추출 (슬라이드 번호 = 페이지)
    - 제목 도형 → heading 블록
    - 본문 도형 → text 블록
    - 표 도형 → table 블록
    - 이미지 내 텍스트 제외 (MVP)
"""
from __future__ import annotations

from uuid import uuid4

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from common_schemas.document import (
    ContentBlock,
    DocumentBlock,
    FileMeta,
    ParserMeta,
    SourceRef,
)

from doc_parser.domain.ports.parser_port import ParserPort


class PptxParser(ParserPort):
    """PowerPoint 파서 구현체.

    지원 MIME 타입:
        application/vnd.openxmlformats-officedocument.presentationml.presentation

    Raises:
        RuntimeError: 파일 손상 또는 읽기 실패 E0202
        RuntimeError: 텍스트 추출 실패 E0203
    """

    MIME_TYPE = (
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation"
    )

    # ──────────────────────────────────────────
    # ParserPort 구현
    # ──────────────────────────────────────────

    def parse(
        self,
        file_path: str,
        file_meta: FileMeta,
    ) -> DocumentBlock:
        """PPTX 파싱 → DocumentBlock 반환.

        Args:
            file_path: PPTX 파일 경로
            file_meta: 파일 메타데이터

        Returns:
            DocumentBlock: 파싱된 문서 블록 (슬라이드별 블록)

        Raises:
            RuntimeError: 파일 손상 (E0202), 텍스트 추출 실패 (E0203)
        """
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise RuntimeError(f"E0202: PPTX 파일 읽기 실패 — {e}") from e

        try:
            blocks: list[ContentBlock] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                block_index = 0

                for shape in slide.shapes:
                    extracted = self._parse_shape(shape, slide_num, block_index)
                    blocks.extend(extracted)
                    block_index += len(extracted)

            return DocumentBlock(
                document_id=uuid4(),
                file_meta=file_meta,
                parser=ParserMeta(
                    parser_name="PptxParser",
                    parser_version="1.0.0",
                ),
                blocks=blocks,
            )

        except Exception as e:
            raise RuntimeError(f"E0203: PPTX 텍스트 추출 실패 — {e}") from e

    def supports(self, mime_type: str) -> bool:
        return mime_type == self.MIME_TYPE

    # ──────────────────────────────────────────
    # Private
    # ──────────────────────────────────────────

    def _parse_shape(
        self,
        shape,
        slide_num: int,
        block_index: int,
    ) -> list[ContentBlock]:
        """도형 타입별 파싱.

        텍스트 도형 → heading / text 블록
        표 도형     → table 블록
        이미지      → MVP 제외
        """
        blocks: list[ContentBlock] = []

        # 표 도형
        if shape.has_table:
            table_data = self._extract_table(shape.table)
            if table_data:
                blocks.append(
                    ContentBlock(
                        block_id=uuid4(),
                        block_type="table",
                        content=None,
                        page=slide_num,
                        table=table_data,
                        source_ref=SourceRef(
                            page=slide_num,
                            block_index=block_index,
                        ),
                    )
                )
            return blocks

        # 텍스트 도형
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                return blocks

            block_type = self._detect_block_type(shape)
            blocks.append(
                ContentBlock(
                    block_id=uuid4(),
                    block_type=block_type,
                    content=text,
                    page=slide_num,
                    source_ref=SourceRef(
                        page=slide_num,
                        block_index=block_index,
                    ),
                )
            )

        return blocks

    def _detect_block_type(self, shape) -> str:
        """도형 이름 또는 placeholder 타입으로 block_type 판단.

        제목 placeholder → heading
        그 외 → text
        """
        # placeholder 타입 확인
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # 1 = TITLE, 13 = CENTER_TITLE
            if ph_type in (1, 13):
                return "heading"
        return "text"

    def _extract_table(self, table) -> list[list[str]]:
        """표 데이터 추출.

        Args:
            table: python-pptx Table 객체

        Returns:
            list[list[str]]: 행 × 열 텍스트 데이터
        """
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        return rows